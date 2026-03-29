#!/usr/bin/env python3
"""生成用户策略PPT，风格参考CDP系统年报"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Emu(12192000)   # 标准16:9
prs.slide_height = Emu(6858000)

# 颜色定义（参考CDP年报风格）
DARK_BLUE = RGBColor(0x06, 0x1F, 0x60)
LIGHT_BLUE = RGBColor(0x00, 0xB0, 0xF0)
LINE_BLUE = RGBColor(0x00, 0x20, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

FONT_CN = '楷体'
FONT_EN = 'Arial'


def add_line(slide, left, top, width, color=LINE_BLUE):
    """添加水平分隔线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(12700))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def set_run(paragraph, text, size=16, bold=False, color=BLACK, font_cn=FONT_CN):
    """设置文本"""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_EN
    run.font.east_asian = font_cn
    return run


def add_content_slide(prs, title_text, sections):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 顶部分隔线
    add_line(slide, Emu(0), Emu(1301843), Emu(12192000))

    # 底部分隔线
    add_line(slide, Emu(0), Emu(6454066), Emu(12192000))

    # 页面标题（蓝色，左上角）
    title_box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(12192000), Emu(1301843))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(12)
    p.space_after = Pt(0)
    # 添加左侧缩进
    set_run(p, '  ' + title_text, size=24, bold=True, color=LIGHT_BLUE)

    # 内容区域
    content_top = Emu(1400000)
    content_left = Emu(223400)
    content_width = Emu(11578000)
    content_height = Emu(4900000)

    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True

    first_section = True
    for section_title, items in sections:
        # 小标题
        if not first_section:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
            set_run(p, '', size=6)

        p = tf.add_paragraph() if not first_section else tf.paragraphs[0]
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        set_run(p, section_title, size=14, bold=True, color=LINE_BLUE)

        # 要点
        for item in items:
            p = tf.add_paragraph()
            p.space_before = Pt(1)
            p.space_after = Pt(1)
            p.level = 0
            # 用bullet符号
            set_run(p, '• ' + item, size=11, bold=False, color=BLACK)

        first_section = False

    return slide


# ==================== 第1页：新客策略 ====================
sections_1 = [
    ('新人券包', [
        '人群：当天注册用户',
        '对照组：1.99（7天）1.99（次日15天）2.99（15天）5折（15天）',
        '实验组：1.99（7天）2.99（次日15天）3.99（15天）5折（15天）',
        '1月19日上线，拉齐后下单+4.8%，杯量+2.5%，单杯实收+2.6%，3日复购率-0.6pp',
        '对首次转化影响较小，主要观察后续复购影响',
        '2月计划：结合档当前实验结果与各方决策新客策略的数据效果',
    ]),
    ('分享有礼', [
        '人群：全量用户群',
        '利益点调整为拉新累计奖励，奖励金额从free drink调整为1.99',
        '用户累计拉新1/2/3/4人分别获得1/3/5/7张1.99券',
        '上线实时push触达策略，推进实物奖励方案',
        '日均拉新36人（完单口径），占大盘拉新6.93%，占比整体在5-7%',
        '1月下旬（25-28日）有明显上升达到9-12%，呈现月末走强趋势',
        '2月计划：1.素材换新，点击率优化  2.实物奖励方案优化',
    ]),
]

add_content_slide(prs, '新客策略进展与计划', sections_1)


# ==================== 第2页：老客策略 ====================
sections_2 = [
    ('活跃用户（近0~15天有交易）', [
        '对照组：6折限品+7折全品 50%  |  实验组：7折全品+7折全品 50%',
        '1月15日上线，拉齐后下单用户+0.5%，杯量+1.2%，单杯实收+0.5%',
        '从0-7天看7折效果最好，建议考虑扩量并拓展为券包形态',
        '2月计划：涨价实验覆盖，对照组6折限品+7折全品，实验组1为7折全品，实验组2为75折全品',
    ]),
    ('浏览未购（近7天浏览未购买）', [
        '对照组：5折券  |  实验组：55折券',
        '1月16日上线，收补后效果显著：下单用户+1.4%，杯量+5.7%，实收+7.6%',
        '3日复购率13.2% vs 9.9%（+3.3pp），用户质量好，建议全量',
        '实时策略1月8日因bug下线，保留周一周五离线触达',
        '2月计划：涨价实验覆盖，对照组5折，实验组1为55折，实验组2为6折',
    ]),
    ('沉默16~30天召回', [
        '对照组：4折券  |  实验组：5折券',
        '1月16日上线，5折提升后单杯实收提升明显，但用户及单量下降，待权衡决策',
        '2月计划：涨价实验覆盖，对照组4折券，实验组1为5折券，实验组2为55折',
    ]),
    ('沉默30天以上召回', [
        '对照组：1.99券包 30%  |  实验组1：3折券包 50%  |  实验组2：4折券包 20%',
        '1月16日上线，拉齐后vs对照组：3折组下单+10.8%实收+6.6%，4折组下单+0.6%实收+17.9%',
        '4折组实收效果最优，但人均杯量-8%，建议观察长期杯量影响后决策',
        '2月计划：涨价实验覆盖，对照组3折券包，实验组1为4折券包，实验组2为6折券包',
    ]),
]

add_content_slide(prs, '老客策略进展与计划', sections_2)


# 保存
output_path = '/Users/xiaoxiao/Vibe coding/用户策略进展与计划.pptx'
prs.save(output_path)
print(f'PPT已保存到: {output_path}')
