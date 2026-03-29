#!/usr/bin/env python3
"""生成用户策略PPT - 看板风格，所有元素可编辑"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xC8, 0x53)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xFF, 0xE8)
SW = 12192000
SH = 6858000


def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s


def hline(slide, l, t, w, c=BLACK, th=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, th)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()


def vline(slide, l, t, h, c=BLACK, th=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, th, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()


def txt(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines = [(text, size, bold, color, align)]"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    for i, (text, sz, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(1); p.space_after = Pt(1)
        r = p.add_run()
        r.text = text; r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = 'Arial'
    return box


E = Emu
L = PP_ALIGN.LEFT
R = PP_ALIGN.RIGHT

# ================================================================
# PAGE 1: 1月核心数据看板
# ================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
rect(s1, 0, 0, SW, SH, fill=WHITE)

# 标题
txt(s1, E(300000), E(80000), E(9000000), E(400000), [
    ('1月核心数据看板 / JAN PERFORMANCE REVIEW', 26, True, BLACK, L)])
txt(s1, E(300000), E(480000), E(9000000), E(280000), [
    ('用户运营效能复盘（User Operation Efficiency Review）', 13, False, GRAY, L)])
txt(s1, E(9500000), E(100000), E(2500000), E(300000), [
    ('● STATUS: COMPLETED', 12, True, GREEN, R)])

# 顶部粗线
hline(s1, E(200000), E(830000), E(SW-400000), BLACK, Pt(3))

# 三列模块标题
cols = [E(250000), E(4150000), E(8050000)]
cw = E(3800000)
mt = E(920000)
mh = E(430000)
titles = [
    ('MODULE A: 获客与转化', 'ACQUISITION & CONVERSION'),
    ('MODULE B: 活跃与留存', 'ACTIVATION & RETENTION'),
    ('MODULE C: 召回与挽留', 'REACTIVATION & CHURN'),
]
for i, (cn, en) in enumerate(titles):
    rect(s1, cols[i], mt, cw, mh, border=BLACK, bw=Pt(2))
    txt(s1, cols[i]+E(80000), mt+E(30000), cw-E(160000), mh-E(60000), [
        (cn, 13, True, BLACK, L), (en, 8, False, GRAY, L)])

# 竖线分隔
ct = E(1350000)
ch = E(5050000)
vline(s1, E(4050000), ct, ch, BLACK, Pt(1.5))
vline(s1, E(7950000), ct, ch, BLACK, Pt(1.5))

# A列中间横线
hline(s1, E(250000), E(3950000), E(3700000), BLACK, Pt(1))
# B列中间横线
hline(s1, E(4150000), E(3700000), E(3700000), BLACK, Pt(1))
hline(s1, E(4150000), E(4800000), E(3700000), BLACK, Pt(1))
# C列中间横线
hline(s1, E(8050000), E(4400000), E(3700000), BLACK, Pt(1))

# ===== MODULE A: 获客与转化 =====

# 新人策略
txt(s1, E(300000), E(1400000), E(3650000), E(2500000), [
    ('新人策略', 20, True, BLACK, L),
    ('EXP: 1.99 vs 2.99', 9, False, GRAY, L),
    ('', 5, False, BLACK, L),
    ('下单 +4.8%', 24, True, GREEN, L),
    ('杯量 +2.5%', 24, True, GREEN, L),
    ('单杯实收 +2.6%', 24, True, GREEN, L),
    ('Revenue/Cup', 8, False, GRAY, L),
    ('', 3, False, BLACK, L),
    ('Insight: 3日复购率微降（-0.6pp），对首次转化', 9, True, BLACK, L),
    ('影响较小，主要观察后续复购影响。', 9, True, BLACK, L),
])

# 分享有礼
txt(s1, E(300000), E(4020000), E(3650000), E(1200000), [
    ('分享有礼', 20, True, BLACK, L),
    ('', 3, False, BLACK, L),
    ('Action: 奖励从 Free Drink 调整为 1.99元', 10, False, BLACK, L),
    ('（有效控制成本）。', 10, False, BLACK, L),
    ('日均拉新36人，占大盘拉新6.93%', 10, False, BLACK, L),
])
# Trend 高亮框
rect(s1, E(310000), E(5350000), E(3550000), E(400000), border=GREEN, bw=Pt(2))
txt(s1, E(360000), E(5370000), E(3450000), E(360000), [
    ('Trend: 1月下旬占比显著上升至 9-12%', 10, True, BLACK, L),
    ('（月末走强）。', 10, True, BLACK, L),
])

# ===== MODULE B: 活跃与留存 =====

# 注册未下单
txt(s1, E(4200000), E(1400000), E(3650000), E(2200000), [
    ('注册未下单', 20, True, BLACK, L),
    ('下单用户数        3日复购率', 11, True, BLACK, L),
    ('+22%      +7.15pp', 28, True, GREEN, L),
    ('', 3, False, BLACK, L),
    ('Tactic: 触达提醒（无额外补券）。', 10, False, BLACK, L),
    ('Status: 用户质量及成本合理，已作为常态化抓手。', 10, False, BLACK, L),
])

# 活跃老客
txt(s1, E(4200000), E(3760000), E(3650000), E(900000), [
    ('活跃老客', 20, True, BLACK, L),
    ('', 3, False, BLACK, L),
    ('Insight: 0-7天内 7折全品 效果最好。', 10, True, BLACK, L),
    ('下单用户+0.5%，杯量+1.2%，单杯实收+0.5%', 10, False, BLACK, L),
])

# 浏览未购
txt(s1, E(4200000), E(4870000), E(3650000), E(1500000), [
    ('浏览未购', 20, True, BLACK, L),
    ('', 3, False, BLACK, L),
    ('Tactic: 5折券收补。', 10, False, BLACK, L),
    ('', 3, False, BLACK, L),
    ('实收 +7.6%', 26, True, GREEN, L),
    ('3日复购率 13.2% (vs 9.9%)', 14, True, BLACK, L),
    ('用户质量好，建议全量。', 10, False, BLACK, L),
])

# ===== MODULE C: 召回与挽留 =====

# 沉默30天+
txt(s1, E(8100000), E(1400000), E(3700000), E(900000), [
    ('沉默30天+', 20, True, BLACK, L),
    ('', 5, False, BLACK, L),
    ('3折组 →  下单 +10.8%', 14, False, BLACK, L),
])
# 4折组高亮框
rect(s1, E(8100000), E(2450000), E(3650000), E(420000), border=GREEN, bw=Pt(2))
txt(s1, E(8160000), E(2470000), E(3550000), E(380000), [
    ('4折组 →  实收 +17.9%', 20, True, BLACK, L),
])
txt(s1, E(8100000), E(2950000), E(3700000), E(1300000), [
    ('', 3, False, BLACK, L),
    ('Winner: 4折组实收效果最优。', 11, True, BLACK, L),
    ('但人均杯量-8%，建议观察长期杯量影响后决策。', 10, False, BLACK, L),
])

# 沉默16-30天
txt(s1, E(8100000), E(4470000), E(3700000), E(1500000), [
    ('沉默16-30天', 20, True, BLACK, L),
    ('', 3, False, BLACK, L),
    ('Result: 5折券提升单杯实收明显，但用户', 10, False, BLACK, L),
    ('及单量下降。待权衡决策。', 10, False, BLACK, L),
])

# 底部
hline(s1, E(200000), E(6450000), E(SW-400000), BLACK, Pt(1.5))
txt(s1, E(300000), E(6480000), E(8000000), E(250000), [
    ('DATA SOURCE: JAN_OPS_LOG_V1 // SYSTEM: OPTIMIZED', 8, False, GRAY, L)])


# ================================================================
# PAGE 2: 2月策略部署
# ================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
rect(s2, 0, 0, SW, SH, fill=WHITE)

# 标题
txt(s2, E(300000), E(80000), E(9000000), E(400000), [
    ('2月策略部署 / FEB STRATEGY DEPLOYMENT', 26, True, BLACK, L)])
txt(s2, E(9500000), E(100000), E(2500000), E(300000), [
    ('■ MODE: EXPERIMENTATION', 12, True, GREEN, R)])

# 核心主线框
rect(s2, E(320000), E(500000), E(8800000), E(340000), border=BLACK, bw=Pt(1.5))
txt(s2, E(400000), E(515000), E(8600000), E(300000), [
    ('核心主线：全链路涨价实验覆盖（Full-link Price Increase Experiments）', 13, False, BLACK, L)])

# 顶部粗线
hline(s2, E(200000), E(920000), E(SW-400000), BLACK, Pt(3))

# 竖线 + 横线
vline(s2, E(6050000), E(1000000), E(5400000), BLACK, Pt(1.5))
hline(s2, E(200000), E(3700000), E(SW-400000), BLACK, Pt(1.5))

# ===== ZONE 1: 新客与增长 =====
txt(s2, E(280000), E(1020000), E(5600000), E(2600000), [
    ('ZONE 1: 新客与增长 (ACQUISITION)', 16, True, BLACK, L),
    ('', 8, False, BLACK, L),
    ('■  素材换新，点击率优化。', 13, False, BLACK, L),
    ('', 5, False, BLACK, L),
    ('■  推进实物奖励方案（随单赠送）。', 13, False, BLACK, L),
    ('', 5, False, BLACK, L),
    ('■  决策新客策略（1.99 / 2.99 / 5折组合）。', 13, False, BLACK, L),
    ('', 5, False, BLACK, L),
    ('■  结合档当前实验结果与各方决策新客策略', 13, False, BLACK, L),
    ('    的数据效果。', 13, False, BLACK, L),
])

# ===== ZONE 2: 活跃用户实验 =====
txt(s2, E(6200000), E(1020000), E(5700000), E(500000), [
    ('ZONE 2: 活跃用户实验 (ACTIVE USERS)', 16, True, BLACK, L),
    ('Frequency: 每天循环发，当天有效', 11, True, BLACK, L),
])

# Split Test
txt(s2, E(6200000), E(1650000), E(1500000), E(1800000), [
    ('', 14, False, BLACK, L),
    ('Split Test', 15, True, BLACK, L),
])

# 分支线
vline(s2, E(7600000), E(1800000), E(1200000), BLACK, Pt(1))
hline(s2, E(7600000), E(1850000), E(300000), BLACK, Pt(1))
hline(s2, E(7600000), E(2350000), E(300000), BLACK, Pt(1))
hline(s2, E(7600000), E(2900000), E(300000), BLACK, Pt(1))

# CONTROL
txt(s2, E(7950000), E(1700000), E(3800000), E(450000), [
    ('CONTROL', 8, False, GRAY, L),
    ('GROUP         6折限品 + 7折全品', 11, False, BLACK, L),
])
# EXP 1
txt(s2, E(7950000), E(2200000), E(3800000), E(450000), [
    ('EXP', 8, False, GRAY, L),
    ('GROUP 1      7折全品 + 7折全品', 11, False, BLACK, L),
])
# EXP 2 高亮
rect(s2, E(7900000), E(2700000), E(3900000), E(450000), fill=LIGHT_GREEN_BG)
txt(s2, E(7950000), E(2720000), E(3800000), E(420000), [
    ('EXP', 8, False, GRAY, L),
    ('GROUP 2      75折全品 + 75折全品', 13, True, BLACK, L),
])

# ===== ZONE 3: 沉睡召回实验 =====
txt(s2, E(280000), E(3780000), E(5600000), E(700000), [
    ('ZONE 3: 沉睡召回实验 (REACTIVATION)', 16, True, BLACK, L),
    ('', 5, False, BLACK, L),
    ('A: Silent 30+ Days', 14, True, BLACK, L),
])
txt(s2, E(280000), E(4450000), E(5600000), E(1000000), [
    ('Cycle: 每周一循环发券包（7日有效）', 10, False, BLACK, L),
    ('CONTROL        3折券包（3折, 2.99, 5折）', 11, False, BLACK, L),
    ('EXP 1            4折券包（4折, 2.99-次日生效, 6折）', 11, False, BLACK, L),
])
# EXP 2 高亮
rect(s2, E(280000), E(5350000), E(5500000), E(300000), fill=LIGHT_GREEN_BG)
txt(s2, E(330000), E(5360000), E(5400000), E(280000), [
    ('EXP 2            6折券包（6折, 3.99-次日生效, 7折）', 12, True, BLACK, L)])

# B: Silent 16-30
hline(s2, E(280000), E(5720000), E(5500000), BLACK, Pt(1))
txt(s2, E(280000), E(5780000), E(5600000), E(600000), [
    ('B: Silent 16-30 Days', 14, True, BLACK, L),
    ('Cycle: 每天循环发，当天有效', 10, False, BLACK, L),
    ('EXP: 测试 55折（vs 4折/5折）', 12, True, BLACK, L),
])

# ===== ZONE 4: 留存防守 =====
txt(s2, E(6200000), E(3780000), E(5700000), E(700000), [
    ('ZONE 4: 留存防守 (RETENTION)', 16, True, BLACK, L),
    ('', 4, False, BLACK, L),
    ('Target: 次月留存（上月消费当月未消费）', 11, True, BLACK, L),
    ('Tactic: 短信触达', 11, True, BLACK, L),
])

# 浏览未购实验
txt(s2, E(6200000), E(4600000), E(2800000), E(300000), [
    ('浏览未购实验', 14, True, BLACK, L),
    ('每天循环发，当天有效', 10, False, BLACK, L),
])

# 实验组列表
txt(s2, E(6200000), E(5000000), E(3200000), E(1300000), [
    ('CONTROL ·················  5折', 12, False, BLACK, L),
    ('', 4, False, BLACK, L),
    ('EXP 1 ·····················  55折', 14, True, BLACK, L),
    ('', 4, False, BLACK, L),
    ('EXP 2 ·····················  6折', 14, True, BLACK, L),
])

# Sensitivity Test
vline(s2, E(9700000), E(5000000), E(1200000), BLACK, Pt(1.5))
txt(s2, E(9900000), E(5200000), E(2000000), E(700000), [
    ('Sensitivity', 16, True, BLACK, L),
    ('Test', 16, True, BLACK, L),
])

# 底部
hline(s2, E(200000), E(6450000), E(SW-400000), BLACK, Pt(1.5))
txt(s2, E(300000), E(6480000), E(9000000), E(250000), [
    ('OBJECTIVE: REVENUE MAXIMIZATION // PRICE ELASTICITY CHECK', 8, False, GRAY, L)])

# 保存
output = '/Users/xiaoxiao/Vibe coding/用户策略进展与计划.pptx'
prs.save(output)
print(f'PPT已保存到: {output}')
