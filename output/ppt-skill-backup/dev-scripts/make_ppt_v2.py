#!/usr/bin/env python3
"""生成用户策略PPT - 参考 Jan Performance / Feb Strategy 看板风格"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Emu(12192000)  # 16:9
prs.slide_height = Emu(6858000)

# 颜色
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xC8, 0x53)  # 强调绿
DARK_GREEN = RGBColor(0x00, 0xA0, 0x40)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x88, 0x88, 0x88)

SLIDE_W = 12192000
SLIDE_H = 6858000


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """添加文本框，lines = [(text, size, bold, color, align), ...]"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor

    for i, line_data in enumerate(lines):
        text, size, bold, color, align = line_data
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Arial'

    return box


def add_hline(slide, left, top, width, color=BLACK, thickness=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_vline(slide, left, top, height, color=BLACK, thickness=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, thickness, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ================================================================
# 第1页：1月核心数据看板
# ================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# 白色背景
add_rect(slide1, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)

# --- 顶部标题区 ---
TITLE_H = Emu(850000)
add_rect(slide1, 0, 0, SLIDE_W, TITLE_H, fill_color=WHITE)

# 主标题
add_text(slide1, Emu(300000), Emu(80000), Emu(9000000), Emu(420000), [
    ('1月核心数据看板 / JAN PERFORMANCE REVIEW', 26, True, BLACK, PP_ALIGN.LEFT),
])
# 副标题
add_text(slide1, Emu(300000), Emu(500000), Emu(9000000), Emu(300000), [
    ('用户运营效能复盘（User Operation Efficiency Review）', 13, False, GRAY, PP_ALIGN.LEFT),
])
# 状态标记
add_text(slide1, Emu(9500000), Emu(100000), Emu(2500000), Emu(350000), [
    ('● STATUS: COMPLETED', 12, True, GREEN, PP_ALIGN.RIGHT),
])

# --- 顶部粗线 ---
add_hline(slide1, Emu(200000), TITLE_H, Emu(SLIDE_W - 400000), BLACK, Pt(3))

# --- 三列模块标题 ---
COL_W = Emu(3800000)
COL_START = Emu(250000)
COL_GAP = Emu(100000)
MOD_TOP = Emu(950000)
MOD_TITLE_H = Emu(450000)

col_positions = [COL_START, Emu(4200000), Emu(8150000)]
col_widths = [Emu(3850000), Emu(3850000), Emu(3850000)]

module_titles = [
    ('MODULE A: 获客与转化', 'ACQUISITION & CONVERSION'),
    ('MODULE B: 活跃与留存', 'ACTIVATION & RETENTION'),
    ('MODULE C: 召回与挽留', 'REACTIVATION & CHURN'),
]

for idx, (title_cn, title_en) in enumerate(module_titles):
    x = col_positions[idx]
    # 模块框
    add_rect(slide1, x, MOD_TOP, col_widths[idx], MOD_TITLE_H, border_color=BLACK, border_width=Pt(2))
    add_text(slide1, x + Emu(80000), MOD_TOP + Emu(40000), col_widths[idx] - Emu(160000), MOD_TITLE_H - Emu(80000), [
        (title_cn, 14, True, BLACK, PP_ALIGN.LEFT),
        (title_en, 9, False, GRAY, PP_ALIGN.LEFT),
    ])

# --- 竖线分隔 ---
CONTENT_TOP = Emu(1400000)
CONTENT_H = Emu(4900000)
add_vline(slide1, Emu(4100000), CONTENT_TOP, CONTENT_H, BLACK, Pt(1.5))
add_vline(slide1, Emu(8050000), CONTENT_TOP, CONTENT_H, BLACK, Pt(1.5))

# --- 横线分隔（A列和B列中间分隔）---
MID_Y = Emu(4000000)
add_hline(slide1, COL_START, MID_Y, Emu(3750000), BLACK, Pt(1))
add_hline(slide1, Emu(4200000), MID_Y, Emu(3750000), BLACK, Pt(1))
add_hline(slide1, Emu(8150000), Emu(4600000), Emu(3750000), BLACK, Pt(1))

# === MODULE A 内容 ===
# 新人策略
add_text(slide1, Emu(300000), Emu(1450000), Emu(3700000), Emu(2500000), [
    ('新人策略', 18, True, BLACK, PP_ALIGN.LEFT),
    ('EXP: 1.99 vs 2.99', 9, False, GRAY, PP_ALIGN.LEFT),
    ('', 6, False, BLACK, PP_ALIGN.LEFT),
    ('下单 +4.8%', 22, True, GREEN, PP_ALIGN.LEFT),
    ('杯量 +2.5%', 22, True, GREEN, PP_ALIGN.LEFT),
    ('单杯实收 +2.6%', 22, True, GREEN, PP_ALIGN.LEFT),
    ('Revenue/Cup', 8, False, GRAY, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('Insight: 3日复购率微降（-0.6pp），对首', 9, True, BLACK, PP_ALIGN.LEFT),
    ('次转化影响较小。', 9, True, BLACK, PP_ALIGN.LEFT),
])

# 分享有礼
add_text(slide1, Emu(300000), Emu(4100000), Emu(3700000), Emu(2200000), [
    ('分享有礼', 18, True, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('Action: 奖励从 Free Drink 调整为 1.99元', 10, False, BLACK, PP_ALIGN.LEFT),
    ('（有效控制成本）。', 10, False, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('日均拉新36人（完单口径），占大盘6.93%', 10, False, BLACK, PP_ALIGN.LEFT),
])
# 绿色边框高亮框
highlight1 = add_rect(slide1, Emu(320000), Emu(5450000), Emu(3500000), Emu(380000),
                       border_color=GREEN, border_width=Pt(2))
add_text(slide1, Emu(370000), Emu(5470000), Emu(3400000), Emu(340000), [
    ('Trend: 1月下旬占比显著上升至 9-12%', 10, True, BLACK, PP_ALIGN.LEFT),
    ('（月末走强）。', 10, True, BLACK, PP_ALIGN.LEFT),
])

# === MODULE B 内容 ===
# 注册未下单
add_text(slide1, Emu(4250000), Emu(1450000), Emu(3700000), Emu(2400000), [
    ('注册未下单', 18, True, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('下单用户数        3日复购率', 12, True, BLACK, PP_ALIGN.LEFT),
    ('+22%      +7.15pp', 26, True, GREEN, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('Tactic: 触达提醒（无额外补券）。Status:', 10, False, BLACK, PP_ALIGN.LEFT),
    ('用户质量及成本合理，已作为常态化抓手。', 10, False, BLACK, PP_ALIGN.LEFT),
    ('', 6, False, BLACK, PP_ALIGN.LEFT),
    ('活跃老客', 16, True, BLACK, PP_ALIGN.LEFT),
    ('Insight: 0-7天内 7折全品 效果最好。', 10, True, BLACK, PP_ALIGN.LEFT),
])

# 浏览未购
add_text(slide1, Emu(4250000), Emu(4100000), Emu(3700000), Emu(2200000), [
    ('浏览未购', 18, True, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('Tactic: 5折券收补。', 10, False, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('实收 +7.6%', 24, True, GREEN, PP_ALIGN.LEFT),
    ('3日复购率 13.2% (vs 9.9%)', 13, True, BLACK, PP_ALIGN.LEFT),
])

# === MODULE C 内容 ===
# 沉默30天+
add_text(slide1, Emu(8200000), Emu(1450000), Emu(3700000), Emu(3000000), [
    ('沉默30天+', 18, True, BLACK, PP_ALIGN.LEFT),
    ('', 6, False, BLACK, PP_ALIGN.LEFT),
    ('3折组 → 下单 +10.8%', 13, False, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
])
# 绿色边框高亮框 - 4折组
highlight2 = add_rect(slide1, Emu(8220000), Emu(2400000), Emu(3600000), Emu(400000),
                       border_color=GREEN, border_width=Pt(2))
add_text(slide1, Emu(8280000), Emu(2420000), Emu(3500000), Emu(360000), [
    ('4折组 → 实收 +17.9%', 18, True, BLACK, PP_ALIGN.LEFT),
])
add_text(slide1, Emu(8200000), Emu(2900000), Emu(3700000), Emu(500000), [
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('Winner: 4折组实收效果最优。', 10, True, BLACK, PP_ALIGN.LEFT),
    ('但人均杯量-8%，建议观察长期影响。', 10, False, BLACK, PP_ALIGN.LEFT),
])

# 沉默16-30天
add_text(slide1, Emu(8200000), Emu(4700000), Emu(3700000), Emu(1500000), [
    ('沉默16-30天', 18, True, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('Result: 5折券提升单杯实收，但用户及单', 10, False, BLACK, PP_ALIGN.LEFT),
    ('量下降。', 10, False, BLACK, PP_ALIGN.LEFT),
])

# --- 底部信息栏 ---
add_hline(slide1, Emu(200000), Emu(6500000), Emu(SLIDE_W - 400000), BLACK, Pt(1.5))
add_text(slide1, Emu(300000), Emu(6530000), Emu(8000000), Emu(250000), [
    ('DATA SOURCE: JAN_OPS_LOG_V1 // SYSTEM: OPTIMIZED', 8, False, GRAY, PP_ALIGN.LEFT),
])


# ================================================================
# 第2页：2月策略部署
# ================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# 白色背景
add_rect(slide2, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)

# --- 顶部标题区 ---
add_text(slide2, Emu(300000), Emu(80000), Emu(9000000), Emu(420000), [
    ('2月策略部署 / FEB STRATEGY DEPLOYMENT', 26, True, BLACK, PP_ALIGN.LEFT),
])
# 状态标记
add_text(slide2, Emu(9500000), Emu(100000), Emu(2500000), Emu(350000), [
    ('■ MODE: EXPERIMENTATION', 12, True, GREEN, PP_ALIGN.RIGHT),
])

# 核心主线（圆角边框效果 - 用矩形模拟）
mainline_box = add_rect(slide2, Emu(320000), Emu(520000), Emu(8500000), Emu(350000),
                         border_color=BLACK, border_width=Pt(1.5))
add_text(slide2, Emu(400000), Emu(540000), Emu(8300000), Emu(300000), [
    ('核心主线：全链路涨价实验覆盖（Full-link Price Increase Experiments）', 13, False, BLACK, PP_ALIGN.LEFT),
])

# --- 顶部粗线 ---
add_hline(slide2, Emu(200000), Emu(950000), Emu(SLIDE_W - 400000), BLACK, Pt(3))

# --- 四区网格 (2x2) ---
GRID_TOP = Emu(1050000)
GRID_MID_Y = Emu(3850000)
GRID_BOT = Emu(6400000)
LEFT_W = Emu(5800000)
RIGHT_W = Emu(5800000)
LEFT_X = Emu(250000)
RIGHT_X = Emu(6200000)

# 竖线
add_vline(slide2, Emu(6100000), GRID_TOP, Emu(5350000), BLACK, Pt(1.5))
# 横线
add_hline(slide2, Emu(200000), GRID_MID_Y, Emu(SLIDE_W - 400000), BLACK, Pt(1.5))

# === ZONE 1: 新客与增长 ===
add_text(slide2, LEFT_X, Emu(1100000), LEFT_W, Emu(2600000), [
    ('ZONE 1: 新客与增长 (ACQUISITION)', 16, True, BLACK, PP_ALIGN.LEFT),
    ('', 8, False, BLACK, PP_ALIGN.LEFT),
    ('■  素材换新，点击率优化。', 12, False, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('■  推进实物奖励方案（随单赠送）。', 12, False, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('■  决策新客策略（1.99 / 2.99 / 5折组合）。', 12, False, BLACK, PP_ALIGN.LEFT),
])

# === ZONE 2: 活跃用户实验 ===
add_text(slide2, RIGHT_X, Emu(1100000), RIGHT_W, Emu(600000), [
    ('ZONE 2: 活跃用户实验 (ACTIVE USERS)', 16, True, BLACK, PP_ALIGN.LEFT),
    ('Frequency: 每天循环发，当天有效', 11, True, BLACK, PP_ALIGN.LEFT),
])

# Split Test 树状图（用文本模拟）
add_text(slide2, RIGHT_X, Emu(1800000), Emu(1500000), Emu(1800000), [
    ('', 10, False, BLACK, PP_ALIGN.LEFT),
    ('', 10, False, BLACK, PP_ALIGN.LEFT),
    ('Split Test', 14, True, BLACK, PP_ALIGN.LEFT),
])

# 分支内容
add_text(slide2, Emu(7800000), Emu(1800000), Emu(4000000), Emu(500000), [
    ('CONTROL       6折限品 + 7折全品', 11, False, BLACK, PP_ALIGN.LEFT),
    ('GROUP', 8, False, GRAY, PP_ALIGN.LEFT),
])
add_text(slide2, Emu(7800000), Emu(2300000), Emu(4000000), Emu(500000), [
    ('EXP            7折全品 + 7折全品', 11, False, BLACK, PP_ALIGN.LEFT),
    ('GROUP 1', 8, False, GRAY, PP_ALIGN.LEFT),
])
# 高亮实验组2
highlight3 = add_rect(slide2, Emu(7750000), Emu(2800000), Emu(4000000), Emu(400000),
                       fill_color=RGBColor(0xE8, 0xFF, 0xE8))
add_text(slide2, Emu(7800000), Emu(2830000), Emu(3900000), Emu(360000), [
    ('EXP            75折全品 + 75折全品', 13, True, BLACK, PP_ALIGN.LEFT),
    ('GROUP 2', 8, False, GRAY, PP_ALIGN.LEFT),
])

# 分支线条
add_hline(slide2, Emu(7500000), Emu(1950000), Emu(250000), BLACK, Pt(1))
add_hline(slide2, Emu(7500000), Emu(2450000), Emu(250000), BLACK, Pt(1))
add_hline(slide2, Emu(7500000), Emu(2950000), Emu(250000), BLACK, Pt(1))
add_vline(slide2, Emu(7500000), Emu(1950000), Emu(1000000), BLACK, Pt(1))

# === ZONE 3: 沉睡召回实验 ===
add_text(slide2, LEFT_X, Emu(3950000), LEFT_W, Emu(2400000), [
    ('ZONE 3: 沉睡召回实验 (REACTIVATION)', 16, True, BLACK, PP_ALIGN.LEFT),
    ('', 6, False, BLACK, PP_ALIGN.LEFT),
    ('A: Silent 30+ Days', 13, True, BLACK, PP_ALIGN.LEFT),
    ('Cycle: 每周一循环发券包（7日有效）', 10, False, BLACK, PP_ALIGN.LEFT),
    ('CONTROL        3折券包', 11, False, BLACK, PP_ALIGN.LEFT),
    ('EXP 1            4折券包（含4折/2.99/6折）', 11, False, BLACK, PP_ALIGN.LEFT),
])
# 高亮 EXP 2
add_rect(slide2, Emu(270000), Emu(5350000), Emu(5500000), Emu(300000),
         fill_color=RGBColor(0xE8, 0xFF, 0xE8))
add_text(slide2, Emu(300000), Emu(5360000), Emu(5400000), Emu(280000), [
    ('EXP 2            6折券包（含6折/3.99/7折）', 11, True, BLACK, PP_ALIGN.LEFT),
])

# 分隔线
add_hline(slide2, LEFT_X, Emu(5750000), Emu(5500000), BLACK, Pt(1))

add_text(slide2, LEFT_X, Emu(5800000), LEFT_W, Emu(600000), [
    ('B: Silent 16-30 Days', 13, True, BLACK, PP_ALIGN.LEFT),
    ('Cycle: 每天循环发', 10, False, BLACK, PP_ALIGN.LEFT),
    ('EXP: 测试 55折（vs 4折/5折）', 11, True, BLACK, PP_ALIGN.LEFT),
])

# === ZONE 4: 留存防守 ===
add_text(slide2, RIGHT_X, Emu(3950000), RIGHT_W, Emu(600000), [
    ('ZONE 4: 留存防守 (RETENTION)', 16, True, BLACK, PP_ALIGN.LEFT),
    ('', 4, False, BLACK, PP_ALIGN.LEFT),
    ('Target: 次月留存（上月消费当月未消费）', 11, True, BLACK, PP_ALIGN.LEFT),
    ('Tactic: 短信触达', 11, True, BLACK, PP_ALIGN.LEFT),
])

# 实验组列表
add_text(slide2, RIGHT_X, Emu(4700000), Emu(3000000), Emu(1400000), [
    ('CONTROL ················· 5折', 12, False, BLACK, PP_ALIGN.LEFT),
    ('', 6, False, BLACK, PP_ALIGN.LEFT),
    ('EXP 1 ····················· 55折', 14, True, BLACK, PP_ALIGN.LEFT),
    ('', 6, False, BLACK, PP_ALIGN.LEFT),
    ('EXP 2 ····················· 6折', 14, True, BLACK, PP_ALIGN.LEFT),
])

# Sensitivity Test 标签
add_text(slide2, Emu(10000000), Emu(4900000), Emu(2000000), Emu(800000), [
    ('Sensitivity', 16, True, BLACK, PP_ALIGN.LEFT),
    ('Test', 16, True, BLACK, PP_ALIGN.LEFT),
])
# 大括号用竖线模拟
add_vline(slide2, Emu(9800000), Emu(4800000), Emu(1200000), BLACK, Pt(1.5))

# --- 底部信息栏 ---
add_hline(slide2, Emu(200000), Emu(6500000), Emu(SLIDE_W - 400000), BLACK, Pt(1.5))
add_text(slide2, Emu(300000), Emu(6530000), Emu(9000000), Emu(250000), [
    ('OBJECTIVE: REVENUE MAXIMIZATION // PRICE ELASTICITY CHECK', 8, False, GRAY, PP_ALIGN.LEFT),
])


# 保存
output_path = '/Users/xiaoxiao/Vibe coding/用户策略进展与计划.pptx'
prs.save(output_path)
print(f'PPT已保存到: {output_path}')
