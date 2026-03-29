#!/usr/bin/env python3
"""
瑞幸美国 PPT 生成器 - LuckinPPT
根据公司模板风格生成标准化 .pptx 文件
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
import os, re, datetime, math

# ============ 样式常量 ============

SLIDE_WIDTH = Emu(24384000)   # 26.67"
SLIDE_HEIGHT = Emu(13716000)  # 15.0"

C = {
    'cover_bg':       RGBColor(0x22, 0x27, 0x73),
    'title':          RGBColor(0x0F, 0x34, 0x60),
    'body':           RGBColor(0x33, 0x33, 0x33),
    'sub':            RGBColor(0x66, 0x66, 0x66),
    'note':           RGBColor(0x99, 0x99, 0x99),
    'gray':           RGBColor(0x88, 0x88, 0x88),
    'red':            RGBColor(0xD3, 0x2F, 0x2F),
    'green':          RGBColor(0x2E, 0x7D, 0x32),
    'table_header':   RGBColor(0x22, 0x27, 0x73),  # 深蓝底
    'table_border':   RGBColor(0xE0, 0xE0, 0xE0),
    'table_stripe':   RGBColor(0xF5, 0xF7, 0xFA),  # 斑马纹
    'conclusion_bg':  RGBColor(0xED, 0xF2, 0xF9),  # 极浅蓝灰
    'white':          RGBColor(0xFF, 0xFF, 0xFF),
    'black':          RGBColor(0x00, 0x00, 0x00),
}

CHART_COLORS = [
    RGBColor(0x22, 0x27, 0x73),  # 深蓝
    RGBColor(0x44, 0x72, 0xC4),  # 中蓝
    RGBColor(0x8B, 0xB8, 0xD0),  # 浅蓝
    RGBColor(0xED, 0x7D, 0x31),  # 橙色
    RGBColor(0xA5, 0xA5, 0xA5),  # 灰色
]

FONT_TITLE = 'STKaiti'
FONT_BODY = 'STKaiti'
ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets')

# 布局常量
CONTENT_LEFT = Emu(1371600)    # 1.5"
CONTENT_WIDTH = Emu(21031200)  # 23"

# 间距常量
GAP_CONCLUSION_SECTION = Emu(274320)  # 0.3" 结论框→章节标题
GAP_SECTION_CONTENT = Emu(91440)      # 0.1" 章节标题→表格/图表
GAP_CONTENT_SECTION = Emu(274320)     # 0.3" 表格/图表→下一章节
GAP_CHART_NEXT = Emu(274320)          # 0.3" 图表→下一元素


# ============ 内部工具 ============

def _parse_colored_text(text):
    """解析颜色标记: <red>text</red>, <green>text</green>, <gray>text</gray>, <bold>text</bold>"""
    parts = []
    pattern = r'<(red|green|gray|bold)>(.*?)</\1>'
    last = 0
    for m in re.finditer(pattern, text):
        if m.start() > last:
            parts.append((text[last:m.start()], None, False))
        tag, content = m.group(1), m.group(2)
        if tag == 'bold':
            parts.append((content, None, True))
        else:
            parts.append((content, C.get(tag), False))
        last = m.end()
    if last < len(text):
        parts.append((text[last:], None, False))
    return parts if parts else [(text, None, False)]


def _add_colored_text(paragraph, text, font_name=None, font_size=Pt(20),
                      default_color=None, default_bold=False):
    """向段落添加带颜色标记的文本"""
    fn = font_name or FONT_BODY
    for txt, color, bold in _parse_colored_text(text):
        run = paragraph.add_run()
        run.text = txt
        run.font.name = fn
        run.font.size = font_size
        run.font.bold = bold or default_bold
        run.font.color.rgb = color if color else (default_color or C['body'])


def _set_para_spacing(paragraph, line_spacing=None, space_before=None, space_after=None):
    """设置段落行间距和段间距"""
    if line_spacing is not None:
        paragraph.line_spacing = line_spacing
    if space_before is not None:
        paragraph.space_before = space_before
    if space_after is not None:
        paragraph.space_after = space_after


def _strip_tags(text):
    """移除颜色标记，返回纯文本"""
    return re.sub(r'<(red|green|gray|bold)>(.*?)</\1>', r'\2', text)


def _estimate_text_height(text, font_size_pt, avail_width_emu, line_spacing=1.4,
                          padding_top=0, padding_bottom=0):
    """
    估算文本渲染后的高度（EMU）。
    - CJK字符宽度 ≈ 1.05 em（含字间距）
    - ASCII字符宽度 ≈ 0.58 em（含字间距）
    - 可用宽度缩减 5% 作为安全边距
    - 最终高度 ×1.2 安全系数（应对 PPT 渲染差异）
    """
    plain = _strip_tags(text)
    font_emu = int(font_size_pt * 12700)  # Pt → EMU
    safe_width = int(avail_width_emu * 0.95)  # 留 5% 安全边距
    # 计算字符宽度总和（EMU）
    total_width = 0
    line_count = 1
    for ch in plain:
        cw = int(font_emu * 1.05) if ord(ch) > 0x7F else int(font_emu * 0.58)
        total_width += cw
        if total_width > safe_width:
            line_count += 1
            total_width = cw  # 当前字符到新行
    line_h = int(font_emu * line_spacing)
    raw_h = int(line_count * line_h + padding_top + padding_bottom)
    return int(raw_h * 1.2)  # 1.2x 安全系数


def _estimate_lines(text, font_size_pt, avail_width_emu):
    """估算文本占多少行"""
    plain = _strip_tags(text)
    font_emu = int(font_size_pt * 12700)
    total_width = 0
    line_count = 1
    for ch in plain:
        cw = font_emu if ord(ch) > 0x7F else int(font_emu * 0.55)
        total_width += cw
        if total_width > avail_width_emu:
            line_count += 1
            total_width = cw
    return line_count


def _no_border(shape):
    """移除形状边框"""
    spPr = shape._element.spPr
    ln = spPr.get_or_add_ln()
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn('a:noFill'))


def _set_cell_border(cell, color=None, width=12700):
    """设置单元格四边边框"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    border_val = str(color or C['table_border'])
    for tag in ['lnL', 'lnR', 'lnT', 'lnB']:
        old = tcPr.find(qn(f'a:{tag}'))
        if old is not None:
            tcPr.remove(old)
        ln = etree.SubElement(tcPr, qn(f'a:{tag}'))
        ln.set('w', str(width))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        clr = etree.SubElement(sf, qn('a:srgbClr'))
        clr.set('val', border_val)


def _set_cell_margins(cell, left=101917, right=101917, top=51117, bottom=51117):
    """设置单元格内边距"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set('marL', str(left))
    tcPr.set('marR', str(right))
    tcPr.set('marT', str(top))
    tcPr.set('marB', str(bottom))


def _style_chart_axes(chart, multi_series=False):
    """统一设置图表坐标轴样式"""
    axis_color = RGBColor(0xCC, 0xCC, 0xCC)
    for axis in [chart.category_axis, chart.value_axis]:
        axis.tick_labels.font.name = FONT_BODY
        axis.tick_labels.font.size = Pt(14)
        axis.format.line.color.rgb = axis_color
    # 网格线：多系列有，单系列无
    chart.value_axis.has_major_gridlines = multi_series
    if multi_series:
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)


def _style_chart_legend(chart, visible=True):
    """设置图例位置为顶部"""
    if visible:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT_BODY
        chart.legend.font.size = Pt(16)
    else:
        chart.has_legend = False


def _add_bar_data_labels(chart):
    """柱状图顶部显示数值"""
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.name = FONT_BODY
    dl.font.size = Pt(14)
    dl.font.color.rgb = C['body']
    dl.number_format_is_linked = False
    try:
        dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass


def _add_pie_data_labels(chart):
    """饼图显示类别+百分比"""
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.name = FONT_BODY
    dl.font.size = Pt(16)
    dl.font.color.rgb = C['body']
    # 通过 XML 设置 showCatName + showPercent
    dLbls = plot._element.get_or_add_dLbls()
    for tag_name, val in [('c:showCatName', '1'), ('c:showPercent', '1'),
                          ('c:showVal', '0'), ('c:showSerName', '0')]:
        for e in dLbls.findall(qn(tag_name)):
            dLbls.remove(e)
        elem = etree.SubElement(dLbls, qn(tag_name))
        elem.set('val', val)


# ============ LuckinPPT 主类 ============

class LuckinPPT:
    def __init__(self, title="报告", author="李宵霄", date=None):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.title = title
        self.author = author
        self.date = date or datetime.date.today().strftime('%Y年%m月%d日')

    def save(self, path):
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)
        print(f"PPT saved: {os.path.abspath(path)}")
        return path

    # ==================== 封面页 ====================

    def add_cover(self, title=None, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = C['cover_bg']
        # 左侧装饰
        p = os.path.join(ASSETS_DIR, 'cover_decoration.png')
        if os.path.exists(p):
            slide.shapes.add_picture(p, Emu(-45720), Emu(1042416), Emu(10810872), Emu(12697464))
        # 右上 Logo
        p = os.path.join(ASSETS_DIR, 'cover_logo.png')
        if os.path.exists(p):
            slide.shapes.add_picture(p, Emu(19553832), Emu(530352), Emu(3840480), Emu(1013064))
        # 主标题 54pt Bold
        tb = slide.shapes.add_textbox(Emu(9601200), Emu(4984272), Emu(12747600), Emu(2926080))
        tf = tb.text_frame
        tf.word_wrap = True
        par = tf.paragraphs[0]
        par.alignment = PP_ALIGN.CENTER
        _set_para_spacing(par, line_spacing=1.2, space_after=Pt(12))
        r = par.add_run()
        r.text = title or self.title
        r.font.name = FONT_TITLE
        r.font.size = Pt(54)
        r.font.bold = True
        r.font.color.rgb = C['white']
        # 副标题 24pt Regular
        sub = subtitle or f"{self.author}\n{self.date}"
        tb2 = slide.shapes.add_textbox(Emu(18025200), Emu(10015200), Emu(2540000), Emu(1524000))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(sub.split('\n')):
            par = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            par.alignment = PP_ALIGN.CENTER
            r = par.add_run()
            r.text = line
            r.font.name = FONT_TITLE
            r.font.size = Pt(24)
            r.font.color.rgb = C['white']
        return slide

    # ==================== 内容页骨架 ====================

    def add_content_page(self, title, subtitle=None):
        """返回 (slide, y_cursor_emu)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # 左上鹿头
        p = os.path.join(ASSETS_DIR, 'logo_topleft.png')
        if os.path.exists(p):
            slide.shapes.add_picture(p, Emu(128016), Emu(0), Emu(1874520), Emu(1893888))
        # 右下 luckin coffee
        p = os.path.join(ASSETS_DIR, 'logo_bottomright.png')
        if os.path.exists(p):
            slide.shapes.add_picture(p, Emu(20574000), Emu(12801600), Emu(3200400), Emu(786384))
        # 大标题 40pt Bold
        tb = slide.shapes.add_textbox(Emu(2167128), Emu(320040), Emu(20878800), Emu(850000))
        tf = tb.text_frame
        tf.word_wrap = True
        par = tf.paragraphs[0]
        _set_para_spacing(par, line_spacing=1.2)
        r = par.add_run()
        r.text = title
        r.font.name = FONT_TITLE
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = C['title']
        # 标题下深蓝色横线（参考PPT标准视觉元素）
        line_y = Emu(1143000)  # 1.25"
        line = slide.shapes.add_shape(1, CONTENT_LEFT, line_y, CONTENT_WIDTH, Emu(27432))  # ~2px
        line.fill.solid()
        line.fill.fore_color.rgb = C['cover_bg']
        _no_border(line)
        y = line_y + Emu(91440)  # 横线下方 0.1"
        # 副标题 20pt #444444（横线下方，清晰可见）
        if subtitle:
            sub_h = _estimate_text_height(subtitle, 20, CONTENT_WIDTH, line_spacing=1.3,
                                          padding_top=45720, padding_bottom=45720)
            sub_h = max(sub_h, 411480)  # 至少 0.45"
            tb2 = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, Emu(sub_h))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            par = tf2.paragraphs[0]
            _set_para_spacing(par, line_spacing=1.3, space_before=Pt(4))
            r = par.add_run()
            r.text = subtitle
            r.font.name = FONT_BODY
            r.font.size = Pt(20)
            r.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
            y += Emu(sub_h)
        return slide, y

    # ==================== 结论框 ====================

    def add_conclusion(self, slide, y, title_text, points):
        """
        双层卡片：深蓝标题栏(白字) + 浅蓝灰内容区
        title_text: 如 "结论：整体可以涨价..."
        points: ["要点1 含<red>+9%</red>标记", ...]
        """
        left = CONTENT_LEFT
        w = CONTENT_WIDTH
        inner_w = w - Emu(365760)  # 左右 padding
        inner_w_val = inner_w  # EMU int value

        # 自适应标题栏高度（28pt 标题）
        title_h = _estimate_text_height(title_text, 28, inner_w_val,
                                        line_spacing=1.3, padding_top=91440, padding_bottom=91440)
        header_h = Emu(max(548640, title_h))  # 至少 0.6"

        # 自适应内容区高度（上下各留 0.15" 内边距，20pt 要点）
        body_pad = Emu(137160)  # 0.15" 上下内边距
        content_h_val = int(body_pad) * 2  # 上下 padding
        for i, pt in enumerate(points):
            full_text = f"{i+1}. {pt}"
            content_h_val += _estimate_text_height(full_text, 20, inner_w_val,
                                                   line_spacing=1.5, padding_top=45720, padding_bottom=45720)
        content_h = Emu(max(548640, content_h_val))  # 至少 0.6"
        box_h = header_h + content_h

        # 深蓝标题栏
        hdr = slide.shapes.add_shape(1, left, y, w, header_h)
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = C['cover_bg']
        _no_border(hdr)
        # 标题栏文字（白色 Bold 28pt）
        tb = slide.shapes.add_textbox(left + Emu(182880), y + Emu(91440),
                                      inner_w, header_h - Emu(182880))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        par = tf.paragraphs[0]
        _add_colored_text(par, title_text, font_size=Pt(28),
                          default_color=C['white'], default_bold=True)

        # 浅蓝灰内容区
        body_y = y + header_h
        body = slide.shapes.add_shape(1, left, body_y, w, content_h)
        body.fill.solid()
        body.fill.fore_color.rgb = C['conclusion_bg']
        _no_border(body)

        # 要点 20pt Regular #333333
        if points:
            tb2 = slide.shapes.add_textbox(left + Emu(182880), body_y + body_pad,
                                           inner_w, content_h - body_pad * 2)
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for i, pt in enumerate(points):
                par = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                _set_para_spacing(par, line_spacing=1.5, space_before=Pt(4), space_after=Pt(4))
                r = par.add_run()
                r.text = f"{i+1}. "
                r.font.name = FONT_BODY
                r.font.size = Pt(20)
                r.font.color.rgb = C['body']
                _add_colored_text(par, pt, font_size=Pt(20), default_color=C['body'])

        return y + box_h + Emu(182880)

    # ==================== 章节标题 ====================

    def add_section_title(self, slide, y, text):
        sec_h = Emu(457200)  # 0.5" 高度
        # 左侧小蓝色竖条装饰
        accent = slide.shapes.add_shape(1, CONTENT_LEFT, y + Emu(68580),
                                        Emu(54864), sec_h - Emu(137160))
        accent.fill.solid()
        accent.fill.fore_color.rgb = C['cover_bg']
        _no_border(accent)
        # 标题文字 26pt
        tb = slide.shapes.add_textbox(CONTENT_LEFT + Emu(137160), y, CONTENT_WIDTH, sec_h)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        par = tf.paragraphs[0]
        r = par.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(26)
        r.font.bold = True
        r.font.color.rgb = C['title']
        return y + sec_h + Emu(45720)

    # ==================== 数据表格 ====================

    def add_table(self, slide, y, headers, rows, col_widths=None, font_size=None,
                  cell_colors=None):
        """
        深蓝表头 + 白色文字 + 斑马纹交替行 + 可选条件着色
        headers: ["列1", "列2"]
        rows: [["值1", "值2 <red>(+5%)</red>"], ...]
        cell_colors: 可选，与rows同shape的2D列表，元素为 RGBColor 或 None
                     非None时覆盖斑马纹背景色
        """
        hdr_size = Pt(20)
        cell_size = font_size or Pt(18)
        nr = len(rows) + 1
        nc = len(headers)
        left = CONTENT_LEFT
        w = CONTENT_WIDTH
        hdr_rh = Emu(548640)   # 表头行高 0.6"
        data_rh = Emu(457200)  # 内容行高 0.5"
        h = hdr_rh + data_rh * len(rows)
        tbl_shape = slide.shapes.add_table(nr, nc, left, y, w, h)
        tbl = tbl_shape.table

        # 列宽
        if col_widths:
            for i, cw in enumerate(col_widths):
                tbl.columns[i].width = cw
        else:
            cw = w // nc
            for i in range(nc):
                tbl.columns[i].width = cw

        # 表头：深蓝底 + 白色文字
        for j, hdr in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = C['table_header']
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            par = cell.text_frame.paragraphs[0]
            par.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            _set_para_spacing(par, line_spacing=1.2, space_before=Pt(4), space_after=Pt(4))
            r = par.add_run()
            r.text = hdr
            r.font.name = FONT_BODY
            r.font.size = hdr_size
            r.font.bold = True
            r.font.color.rgb = C['white']
            _set_cell_border(cell)
            _set_cell_margins(cell)

        # 数据行：斑马纹（偶数行 #F5F7FA，奇数行白色）+ 条件着色
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i + 1, j)
                cell.text = ""
                cell.fill.solid()
                # 条件着色优先，否则斑马纹
                custom_bg = cell_colors[i][j] if cell_colors and cell_colors[i][j] else None
                cell.fill.fore_color.rgb = custom_bg if custom_bg else (
                    C['table_stripe'] if i % 2 == 1 else C['white'])
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                par = cell.text_frame.paragraphs[0]
                par.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                _set_para_spacing(par, line_spacing=1.2, space_before=Pt(4), space_after=Pt(4))
                _add_colored_text(par, str(val), font_size=cell_size, default_color=C['body'])
                _set_cell_border(cell)
                _set_cell_margins(cell)

        return y + h + Emu(137160)  # 表格后间距收紧到 0.15"

    # ==================== 图表 ====================

    def _style_chart_title(self, chart, title):
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = title
            for r in chart.chart_title.text_frame.paragraphs[0].runs:
                r.font.name = FONT_BODY
                r.font.size = Pt(22)
                r.font.bold = True
                r.font.color.rgb = C['body']

    def add_bar_chart(self, slide, y, title, categories, series_dict,
                      width=None, height=None, stacked=False):
        """柱状图：柱顶数据标签，间隙比 100%"""
        cd = CategoryChartData()
        cd.categories = categories
        for name, vals in series_dict.items():
            cd.add_series(name, vals)
        left = CONTENT_LEFT
        w = width or CONTENT_WIDTH
        h = height or Emu(5486400)  # 6" 默认高度
        ct = XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED
        cf = slide.shapes.add_chart(ct, left, y, w, h, cd)
        chart = cf.chart
        multi = len(series_dict) > 1
        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]
        _style_chart_legend(chart, visible=multi)
        _style_chart_axes(chart, multi_series=multi)
        chart.plots[0].gap_width = 100
        _add_bar_data_labels(chart)
        self._style_chart_title(chart, title)
        return y + h + GAP_CHART_NEXT

    def add_line_chart(self, slide, y, title, categories, series_dict,
                       width=None, height=None):
        """折线图：2.5pt 线宽，圆形标记点 8pt，不显示数据标签"""
        cd = CategoryChartData()
        cd.categories = categories
        for name, vals in series_dict.items():
            cd.add_series(name, vals)
        left = CONTENT_LEFT
        w = width or CONTENT_WIDTH
        h = height or Emu(5486400)  # 6" 默认高度
        cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, y, w, h, cd)
        chart = cf.chart
        multi = len(series_dict) > 1
        for i, s in enumerate(chart.series):
            color = CHART_COLORS[i % len(CHART_COLORS)]
            s.format.line.color.rgb = color
            s.format.line.width = Pt(2.5)
            # 圆形标记点 8pt
            ser = s._element
            marker = ser.find(qn('c:marker'))
            if marker is None:
                marker = etree.SubElement(ser, qn('c:marker'))
            for tag, val in [('c:symbol', 'circle'), ('c:size', '8')]:
                el = marker.find(qn(tag))
                if el is None:
                    el = etree.SubElement(marker, qn(tag))
                el.set('val', val)
            # 标记点填充色
            spPr = marker.find(qn('c:spPr'))
            if spPr is None:
                spPr = etree.SubElement(marker, qn('c:spPr'))
            else:
                for child in list(spPr):
                    spPr.remove(child)
            sf = etree.SubElement(spPr, qn('a:solidFill'))
            clr = etree.SubElement(sf, qn('a:srgbClr'))
            clr.set('val', str(color))
        _style_chart_legend(chart, visible=multi)
        _style_chart_axes(chart, multi_series=multi)
        self._style_chart_title(chart, title)
        return y + h + GAP_CHART_NEXT

    def add_pie_chart(self, slide, y, title, categories, values,
                      width=None, height=None):
        """饼图：类别+百分比标签"""
        cd = CategoryChartData()
        cd.categories = categories
        cd.add_series('', values)
        left = CONTENT_LEFT
        w = width or Emu(10637600)
        h = height or Emu(5486400)  # 6" 默认高度
        cf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, y, w, h, cd)
        chart = cf.chart
        _style_chart_legend(chart, visible=True)
        series = chart.series[0]
        for i in range(len(categories)):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]
        _add_pie_data_labels(chart)
        self._style_chart_title(chart, title)
        return y + h + GAP_CHART_NEXT

    # ==================== 文本块 ====================

    def add_text_block(self, slide, y, text, font_size=Pt(20), color=None, bold=False):
        # 自适应高度估算
        fs_pt = font_size / 12700  # Pt(x) returns EMU int, convert back to pt
        total_h = 0
        for line in text.split('\n'):
            total_h += _estimate_text_height(line, fs_pt, CONTENT_WIDTH,
                                             line_spacing=1.5, padding_top=50800, padding_bottom=50800)
        box_h = Emu(max(274320, total_h))
        tb = slide.shapes.add_textbox(CONTENT_LEFT, y, CONTENT_WIDTH, box_h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split('\n')):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _set_para_spacing(par, line_spacing=1.5, space_before=Pt(4), space_after=Pt(4))
            _add_colored_text(par, line, font_size=font_size,
                              default_color=color or C['body'], default_bold=bold)
        return y + box_h + Emu(91440)

    # ==================== 脚注 ====================

    def add_footnote(self, slide, text, y=None):
        fy = y or Emu(13262616)
        tb = slide.shapes.add_textbox(CONTENT_LEFT, fy, CONTENT_WIDTH, Emu(320040))
        tf = tb.text_frame
        tf.word_wrap = True
        par = tf.paragraphs[0]
        _set_para_spacing(par, line_spacing=1.2)
        r = par.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(14)
        r.font.color.rgb = C['note']
        return fy + Emu(320040)


# ============ 快捷函数 ============

def create_monthly_report(title, author, date, pages_data, output_path):
    """
    pages_data: [
        {
            "title": "页面标题",
            "subtitle": "可选副标题",
            "conclusion": {"title": "结论：xxx", "points": ["要点1", ...]},
            "sections": [
                {"title": "章节标题", "table": {"headers": [...], "rows": [...]}},
                {"title": "章节标题", "chart": {"type": "bar|line|pie", ...}},
                {"text": "普通文本段落"},
            ],
            "footnote": "脚注"
        }, ...
    ]
    """
    ppt = LuckinPPT(title=title, author=author, date=date)
    ppt.add_cover()
    for page in pages_data:
        slide, y = ppt.add_content_page(page['title'], page.get('subtitle'))
        if 'conclusion' in page:
            y = ppt.add_conclusion(slide, y,
                                   page['conclusion']['title'],
                                   page['conclusion'].get('points', []))
        for sec in page.get('sections', []):
            if 'title' in sec:
                y = ppt.add_section_title(slide, y, sec['title'])
            if 'table' in sec:
                t = sec['table']
                y = ppt.add_table(slide, y, t['headers'], t['rows'],
                                  t.get('col_widths'), t.get('font_size'),
                                  t.get('cell_colors'))
            if 'chart' in sec:
                c = sec['chart']
                if c['type'] == 'bar':
                    y = ppt.add_bar_chart(slide, y, c.get('title', ''),
                                          c['categories'], c['series'],
                                          stacked=c.get('stacked', False))
                elif c['type'] == 'line':
                    y = ppt.add_line_chart(slide, y, c.get('title', ''),
                                           c['categories'], c['series'])
                elif c['type'] == 'pie':
                    y = ppt.add_pie_chart(slide, y, c.get('title', ''),
                                          c['categories'], c['values'])
            if 'text' in sec:
                y = ppt.add_text_block(slide, y, sec['text'])
        if 'footnote' in page:
            ppt.add_footnote(slide, page['footnote'])
    return ppt.save(output_path)
